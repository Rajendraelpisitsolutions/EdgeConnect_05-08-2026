"""
File:     docs/marketing/assets/build-pitch-deck-v1.py
Purpose:  Generate the executive pitch deck (pitch-deck-v1.pptx) for the
          Elpis Industrial Intelligence Platform.
Source:   docs/marketing/pitch-deck-outline-v1.md (12 slides, locked)
Tokens:   docs/marketing/assets/brand/BRAND_TOKENS.md v1 (locked)
Spec:     dark premium-industrial, single teal accent, Inter font,
          left vertical accent rule motif on every content slide
Output:   docs/marketing/assets/pitch-deck-v1.pptx
Run:      python docs/marketing/assets/build-pitch-deck-v1.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Brand tokens (mirrors BRAND_TOKENS.md v1)
# ---------------------------------------------------------------------------

BG_DEFAULT       = RGBColor(0x1A, 0x1F, 0x26)
BG_DEEP          = RGBColor(0x0F, 0x14, 0x19)
SURFACE_HERO     = RGBColor(0x2A, 0x2F, 0x36)
SURFACE_HERO_BR  = RGBColor(0x23, 0x28, 0x30)   # bottom-right of hero gradient
SURFACE_SECOND   = RGBColor(0x3A, 0x40, 0x49)
BORDER_SUBTLE    = RGBColor(0x4A, 0x55, 0x60)
BORDER_STRONG    = RGBColor(0x5E, 0x6B, 0x78)
TEXT_BODY        = RGBColor(0xE8, 0xEC, 0xF1)
TEXT_MUTED       = RGBColor(0xA8, 0xB3, 0xBD)
TEXT_HEADING     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_CAPTION     = RGBColor(0xC8, 0xD0, 0xD8)
BRAND_TEAL       = RGBColor(0x00, 0xA0, 0xE0)

FONT_PRIMARY = "Inter"   # falls back to Calibri/system on machines without Inter

# Slide dimensions (LAYOUT_WIDE 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout helpers
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.7)
MARGIN_T = Inches(0.55)
MARGIN_B = Inches(0.5)
CONTENT_W = Inches(13.333 - 1.4)  # 11.93"
ACCENT_RULE_W = Inches(0.08)

ASSETS_DIR = Path(__file__).parent
OUT_PPTX = ASSETS_DIR / "pitch-deck-v1.pptx"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill_rgb is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb
        if line_w is not None:
            s.line.width = line_w
    # Remove default text body padding
    tf = s.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return s


def add_text(slide, text, x, y, w, h, *,
             font=FONT_PRIMARY, size=14, bold=False, italic=False,
             color=TEXT_BODY, align="left", anchor="top",
             letter_spacing=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        # letter spacing in PPTX is set via OOXML "spc" attribute (hundredths of a point)
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(letter_spacing * 100)))
    return tb


def add_multiline(slide, paragraphs, x, y, w, h, *,
                  align="left", anchor="top", line_spacing=1.3,
                  para_space_after=0):
    """
    paragraphs: list of paragraph dicts. EACH entry is its own paragraph.
      {"text": str, "size": pt, "bold": bool, "italic": bool, "color": RGBColor,
       "font": str, "space_after": pt, "space_before": pt, "align": str}
    """
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM
    }[anchor]

    for i, r in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[r.get("align", align)]
        p.line_spacing = line_spacing
        if r.get("space_before"):
            p.space_before = Pt(r["space_before"])
        if r.get("space_after") is not None:
            p.space_after = Pt(r["space_after"])
        elif para_space_after:
            p.space_after = Pt(para_space_after)
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", FONT_PRIMARY)
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", TEXT_BODY)
    return tb


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def speaker_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# Visual motif: left accent rule on every content slide
def content_chrome(slide, slide_num, total=12):
    set_bg(slide, BG_DEFAULT)
    # Left teal accent rule
    add_rect(slide, Inches(0), Inches(0), ACCENT_RULE_W, SLIDE_H, fill_rgb=BRAND_TEAL)
    # Footer: brand wordmark text on left, slide number on right
    add_text(slide, "ELPIS  ·  INDUSTRIAL INTELLIGENCE PLATFORM",
             Inches(0.4), Inches(7.05), Inches(7), Inches(0.3),
             size=9, color=TEXT_MUTED, letter_spacing=2.2)
    add_text(slide, f"{slide_num:02d} / {total:02d}",
             Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
             size=9, color=TEXT_MUTED, align="right", letter_spacing=1.5)


def title_chrome(slide):
    """Title and closing slides — no left accent rule, deeper bg."""
    set_bg(slide, BG_DEEP)


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# SLIDE 1 — Title
# ============================================================
s = add_blank(prs)
title_chrome(s)

# Subtle teal accent line above the title (a small visual signature)
add_rect(s, Inches(0.95), Inches(2.55), Inches(0.6), Inches(0.04), fill_rgb=BRAND_TEAL)

# Pre-title small label
add_text(s, "ELPIS IT SOLUTIONS  ·  2026",
         Inches(0.95), Inches(2.20), Inches(8), Inches(0.35),
         size=12, color=TEXT_MUTED, letter_spacing=3.0)

# Title
add_text(s, "Industrial Intelligence Platform",
         Inches(0.95), Inches(2.75), Inches(11.5), Inches(1.4),
         size=60, bold=True, color=TEXT_HEADING, line_spacing=1.05)

# Subtitle
add_text(s, "Unified industrial connectivity and operational intelligence for modern manufacturing.",
         Inches(0.95), Inches(4.2), Inches(11), Inches(1.0),
         size=22, color=TEXT_MUTED, italic=True, line_spacing=1.3)

# Presenter block bottom-left
add_multiline(s, [
    {"text": "Presented by",  "size": 11, "color": TEXT_MUTED},
    {"text": "[ Presenter Name ]",  "size": 16, "bold": True, "color": TEXT_BODY},
    {"text": "[ Role ]  ·  [ Date ]",  "size": 12, "color": TEXT_MUTED},
], Inches(0.95), Inches(6.1), Inches(6), Inches(0.9), line_spacing=1.4)

# Elpis wordmark placeholder bottom-right
add_text(s, "ELPIS", Inches(11.4), Inches(6.4), Inches(1.5), Inches(0.5),
         size=24, bold=True, color=BRAND_TEAL, align="right", letter_spacing=4.0)
add_text(s, "THINK · CREATE · ENABLE", Inches(10.0), Inches(6.8), Inches(2.9), Inches(0.3),
         size=8, color=TEXT_MUTED, align="right", letter_spacing=2.0)

speaker_notes(s,
    "Don't read the slide. Open with one sentence: 'I'm going to show you how to put your "
    "whole plant — every machine, every shift, every controller — on one operational view, "
    "without ripping out anything you already own.'")

# ============================================================
# SLIDE 2 — The problem
# ============================================================
s = add_blank(prs)
content_chrome(s, 2)

# Section label
add_text(s, "THE PROBLEM",
         Inches(0.7), Inches(0.5), Inches(6), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

# Headline
add_text(s, "The data is already on the floor.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=40, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# Two-column layout: left=bullets, right=punchline
# Left column — three problem rows with leading teal dot
col_left_x = Inches(0.7)
col_right_x = Inches(7.2)

row_ys = [Inches(2.55), Inches(3.85), Inches(5.15)]
problems = [
    ("Fanuc, Brother, Siemens, energy meters",
     "Each speaks a different language."),
    ("OEE numbers stitched together from spreadsheets",
     "And from operator memory at the end of the shift."),
    ("Downtime detected in hindsight",
     "Not in the moment it happens."),
]
for y, (lead, follow) in zip(row_ys, problems):
    # Teal dot
    add_rect(s, col_left_x, y + Inches(0.18), Inches(0.18), Inches(0.18),
             fill_rgb=BRAND_TEAL, shape=MSO_SHAPE.OVAL)
    # Lead text
    add_text(s, lead, col_left_x + Inches(0.4), y, Inches(5.8), Inches(0.45),
             size=18, bold=True, color=TEXT_BODY, line_spacing=1.2)
    # Follow text
    add_text(s, follow, col_left_x + Inches(0.4), y + Inches(0.5), Inches(5.8), Inches(0.45),
             size=15, color=TEXT_MUTED, italic=False, line_spacing=1.3)

# Vertical divider line between columns (subtle)
add_rect(s, Inches(6.95), Inches(2.4), Inches(0.02), Inches(3.5),
         fill_rgb=BORDER_SUBTLE)

# Right column — punchline
add_multiline(s, [
    {"text": "Plants don't have a", "size": 30, "color": TEXT_MUTED, "italic": True},
    {"text": "data problem.", "size": 34, "bold": True, "color": TEXT_HEADING},
    {"text": "They have a", "size": 30, "color": TEXT_MUTED, "italic": True},
    {"text": "decision problem.", "size": 34, "bold": True, "color": BRAND_TEAL},
], col_right_x, Inches(2.55), Inches(5.4), Inches(3.6),
   line_spacing=1.2, para_space_after=4)

speaker_notes(s,
    "Anchor with a specific scenario: 'Your line manager finds out about a 47-minute "
    "downtime event at 6 AM the next morning, when she opens her email.' Pause. Let it land.")

# ============================================================
# SLIDE 3 — Designed for
# ============================================================
s = add_blank(prs)
content_chrome(s, 3)

add_text(s, "DESIGNED FOR",
         Inches(0.7), Inches(0.5), Inches(6), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Five places this platform was built for.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=36, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# Five rows: number + label
audiences = [
    ("01", "Multi-vendor CNC manufacturing plants",
     "Mixed Fanuc, Brother, Mazak fleets on a single operational view."),
    ("02", "Automotive parts and precision machining operations",
     "Tier-2 suppliers needing audit-defensible OEE and provenance."),
    ("03", "Brownfield modernization projects",
     "Older controllers brought into a modern analytics stack — the iron stays."),
    ("04", "OEM machine monitoring deployments",
     "Trust-respecting telemetry on machines deployed at customer sites."),
    ("05", "Multi-site industrial operations teams",
     "Fleet coherence across plants with one EREMOS V2 tenant."),
]
start_y = 2.05
row_h = 0.95
for i, (num, lead, sub) in enumerate(audiences):
    y = Inches(start_y + i * row_h)
    # Number in teal
    add_text(s, num, Inches(0.7), y, Inches(1.0), Inches(0.7),
             size=30, bold=True, color=BRAND_TEAL, anchor="middle")
    # Lead label
    add_text(s, lead, Inches(1.85), y, Inches(7.0), Inches(0.4),
             size=18, bold=True, color=TEXT_BODY)
    # Sub description
    add_text(s, sub, Inches(1.85), y + Inches(0.42), Inches(10.5), Inches(0.35),
             size=13, color=TEXT_MUTED, italic=False)
    # Thin separator line under each row (except last)
    if i < len(audiences) - 1:
        add_rect(s, Inches(1.85), y + Inches(0.85), Inches(10.5), Inches(0.012),
                 fill_rgb=BORDER_SUBTLE)

speaker_notes(s,
    "Pause at each bullet. Pick the one your prospect is, read it directly to them, then "
    "move on. If the prospect doesn't fit any of these five, this is the moment to qualify "
    "out gracefully.")

# ============================================================
# SLIDE 4 — The solution
# ============================================================
s = add_blank(prs)
content_chrome(s, 4)

add_text(s, "THE SOLUTION",
         Inches(0.7), Inches(0.5), Inches(6), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Two products. One platform. One operational view.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=34, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# Two hero blocks: EdgeConnect (left) and EREMOS V2 (right) with arrow in between
hero_y = Inches(2.4)
hero_h = Inches(3.6)
hero_w = Inches(4.8)
left_x = Inches(0.7)
right_x = Inches(7.85)

# Left hero — EdgeConnect
add_rect(s, left_x, hero_y, hero_w, hero_h, fill_rgb=SURFACE_HERO, line_rgb=BORDER_STRONG)
# Small teal accent corner
add_rect(s, left_x + Inches(0.35), hero_y + Inches(0.35), Inches(0.4), Inches(0.06),
         fill_rgb=BRAND_TEAL)
add_text(s, "EdgeConnect",
         left_x + Inches(0.35), hero_y + Inches(0.55), Inches(4.2), Inches(0.7),
         size=34, bold=True, color=TEXT_HEADING)
add_text(s, "Edge runtime",
         left_x + Inches(0.35), hero_y + Inches(1.2), Inches(4.2), Inches(0.4),
         size=15, color=TEXT_MUTED, italic=True)
add_rect(s, left_x + Inches(0.35), hero_y + Inches(1.7), Inches(4.1), Inches(0.012),
         fill_rgb=BORDER_SUBTLE)
add_multiline(s, [
    {"text": "Collects from every controller on your floor.",
     "size": 14, "color": TEXT_BODY},
    {"text": " ", "size": 6},
    {"text": "Canonical data pipeline. Per-route store-and-forward. Three-way diagnostics. "
             "Signed offline licensing. Deploys per plant — one per site.",
     "size": 13, "color": TEXT_MUTED, "italic": True},
], left_x + Inches(0.35), hero_y + Inches(1.85), Inches(4.1), Inches(1.6),
   line_spacing=1.4)

# Right hero — EREMOS V2
add_rect(s, right_x, hero_y, hero_w, hero_h, fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE)
add_rect(s, right_x + Inches(0.35), hero_y + Inches(0.35), Inches(0.4), Inches(0.06),
         fill_rgb=BRAND_TEAL)
add_text(s, "EREMOS V2",
         right_x + Inches(0.35), hero_y + Inches(0.55), Inches(4.2), Inches(0.7),
         size=34, bold=True, color=TEXT_HEADING)
add_text(s, "Industrial intelligence",
         right_x + Inches(0.35), hero_y + Inches(1.2), Inches(4.2), Inches(0.4),
         size=15, color=TEXT_MUTED, italic=True)
add_rect(s, right_x + Inches(0.35), hero_y + Inches(1.7), Inches(4.1), Inches(0.012),
         fill_rgb=BORDER_SUBTLE)
add_multiline(s, [
    {"text": "Turns that data into OEE, alarms, incidents, and reports.",
     "size": 14, "color": TEXT_BODY},
    {"text": " ", "size": 6},
    {"text": "Multi-tenant analytics. OEE via Segments. Persistent alarm + incident workflow. "
             "Configurable alerting. PDF / Excel reports. Tool-life · tag mapping.",
     "size": 13, "color": TEXT_MUTED, "italic": True},
], right_x + Inches(0.35), hero_y + Inches(1.85), Inches(4.1), Inches(1.6),
   line_spacing=1.4)

# Arrow + label between them (centered horizontally between heroes)
arrow_y = hero_y + Inches(1.65)
arrow_x_start = left_x + hero_w + Inches(0.05)   # 5.55
arrow_x_end   = right_x - Inches(0.05)            # 7.80
arrow_mid_x   = (arrow_x_start + arrow_x_end) / 2

# Horizontal arrow (filled triangle on right)
add_rect(s, arrow_x_start, arrow_y + Inches(0.16), arrow_x_end - arrow_x_start, Inches(0.06),
         fill_rgb=BRAND_TEAL)
add_rect(s, arrow_x_end - Inches(0.18), arrow_y + Inches(0.08), Inches(0.22), Inches(0.22),
         fill_rgb=BRAND_TEAL, shape=MSO_SHAPE.RIGHT_TRIANGLE)
# Actually use ISOSCELES_TRIANGLE rotated 90 — simpler: just use RIGHT_TRIANGLE oriented right
# label
add_text(s, "MQTT  ·  OPC UA",
         arrow_x_start, arrow_y - Inches(0.45), arrow_x_end - arrow_x_start, Inches(0.4),
         size=11, color=BRAND_TEAL, align="center", letter_spacing=3.0, bold=True)

# Tagline at bottom centered
add_text(s, "From the spindle to the dashboard, on one foundation.",
         Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5),
         size=18, color=TEXT_CAPTION, italic=True, align="center")

speaker_notes(s,
    "Establish that this is two products that work together, sold and licensed independently. "
    "Don't dive into how yet — that's later.")

# ============================================================
# SLIDE 5 — Outcomes you can hold us to
# ============================================================
s = add_blank(prs)
content_chrome(s, 5)

add_text(s, "OUTCOMES YOU CAN HOLD US TO",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Six concrete outcomes. Operational deliverables — not promises.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=30, bold=True, color=TEXT_HEADING, line_spacing=1.1)

outcomes = [
    ("Cut unplanned downtime",
     "Persistent alarm tracking and incident workflows."),
    ("Trust your OEE number",
     "Every input collected at the controller."),
    ("Modernize legacy controllers",
     "No replacements required."),
    ("See your whole fleet",
     "Multiple plants, one operational view."),
    ("Keep sensitive data where it belongs",
     "Fully offline-capable at the edge."),
    ("Pass your audit",
     "Hash-chained config history, signed offline licensing."),
]

# 2-column x 3-row grid
card_w = Inches(5.9)
card_h = Inches(1.35)
col_xs = [Inches(0.7), Inches(6.75)]
row_ys = [Inches(2.4), Inches(3.95), Inches(5.5)]

for i, (lead, sub) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    x = col_xs[col]
    y = row_ys[row]
    # Card
    add_rect(s, x, y, card_w, card_h,
             fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE)
    # Teal indicator
    add_rect(s, x, y, Inches(0.06), card_h, fill_rgb=BRAND_TEAL)
    # Lead
    add_text(s, lead, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.55), Inches(0.5),
             size=18, bold=True, color=TEXT_HEADING)
    # Sub
    add_text(s, sub, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.55), Inches(0.5),
             size=13, color=TEXT_MUTED)

speaker_notes(s,
    "This is the slide to spend time on. Read each outcome out loud. Ask the prospect "
    "which one matters most to them right now. The answer shapes the rest of the conversation.")

# ============================================================
# SLIDE 6 — Replace spreadsheet operations
# ============================================================
s = add_blank(prs)
content_chrome(s, 6)

add_text(s, "REPLACE SPREADSHEET OPERATIONS",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Most plants already have the data.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.7),
         size=34, bold=True, color=TEXT_HEADING, line_spacing=1.1)
add_text(s, "What they lack is a system that produces:",
         Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.55),
         size=22, color=TEXT_MUTED, italic=True, line_spacing=1.2)

# Before/after split
# Left: faded "spreadsheet" panel (the before)
before_x = Inches(0.7)
before_w = Inches(5.0)
before_y = Inches(2.7)
before_h = Inches(3.5)

add_rect(s, before_x, before_y, before_w, before_h,
         fill_rgb=SURFACE_SECOND, line_rgb=BORDER_SUBTLE)
# Spreadsheet mock
add_text(s, "TODAY", before_x + Inches(0.3), before_y + Inches(0.25), Inches(3), Inches(0.3),
         size=11, color=TEXT_MUTED, letter_spacing=2.5, bold=True)
# fake spreadsheet rows
mock_y = before_y + Inches(0.75)
for i in range(8):
    row_y = mock_y + Inches(0.32 * i)
    # row line
    add_rect(s, before_x + Inches(0.3), row_y, before_w - Inches(0.6), Inches(0.018),
             fill_rgb=BORDER_SUBTLE)
    # fake cells
    for c in range(4):
        cell_x = before_x + Inches(0.3 + c * 1.05)
        add_text(s, "—————", cell_x, row_y + Inches(0.05), Inches(1.0), Inches(0.25),
                 size=10, color=TEXT_MUTED if i % 2 == 0 else BORDER_SUBTLE, line_spacing=1.0)
add_text(s, "Tribal knowledge in cells.\nNo timestamps you can trust.",
         before_x + Inches(0.3), before_y + before_h - Inches(0.75),
         before_w - Inches(0.6), Inches(0.65),
         size=12, color=TEXT_MUTED, italic=True, line_spacing=1.3)

# Right: the five outcomes
after_x = Inches(6.4)
after_w = Inches(6.2)
after_y = Inches(2.7)

add_text(s, "WITH THE PLATFORM",
         after_x, after_y, Inches(6), Inches(0.3),
         size=11, color=BRAND_TEAL, letter_spacing=2.5, bold=True)

replacements = [
    ("Trusted timestamps",  "Captured at the controller — not transcribed from a clipboard."),
    ("Auditable OEE",       "Segment-based math you can show an auditor."),
    ("Persistent alarm history", "Every fault on the record."),
    ("Unified machine visibility", "One view across CNCs, PLCs, meters."),
    ("Centralized workflows", "Shift reports as a record — not a phone call."),
]
list_y = after_y + Inches(0.55)
for i, (lead, sub) in enumerate(replacements):
    y = list_y + Inches(i * 0.66)
    # bullet square
    add_rect(s, after_x, y + Inches(0.16), Inches(0.14), Inches(0.14), fill_rgb=BRAND_TEAL)
    add_text(s, lead, after_x + Inches(0.32), y, Inches(5.5), Inches(0.32),
             size=15, bold=True, color=TEXT_BODY)
    add_text(s, sub, after_x + Inches(0.32), y + Inches(0.3), Inches(5.5), Inches(0.32),
             size=12, color=TEXT_MUTED)

speaker_notes(s,
    "Ask the prospect: 'How many spreadsheets does your shift handover currently depend on?' "
    "Most prospects say three to seven. That's the opening.")

# ============================================================
# SLIDE 7 — Connectivity coverage
# ============================================================
s = add_blank(prs)
content_chrome(s, 7)

add_text(s, "CONNECTIVITY COVERAGE",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "We support the protocols your controllers already speak.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=30, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# 2x2 grid of protocol category cards
quad_w = Inches(5.9)
quad_h = Inches(1.8)
quad_xs = [Inches(0.7), Inches(6.75)]
quad_ys = [Inches(2.3), Inches(4.25)]

categories = [
    ("SOUTHBOUND", "CNC controllers",
     "FOCAS2  ·  MT-LINKi  ·  MTConnect  ·  Brother HTTP"),
    ("SOUTHBOUND", "PLC + instrumentation",
     "Modbus TCP"),
    ("NORTHBOUND", "Messaging",
     "MQTT (any compliant broker)"),
    ("NORTHBOUND", "Enterprise integration",
     "OPC UA Server"),
]

for i, (direction, title, protocols) in enumerate(categories):
    col = i % 2; row = i // 2
    x = quad_xs[col]; y = quad_ys[row]
    add_rect(s, x, y, quad_w, quad_h, fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE)
    # Direction label
    add_text(s, direction, x + Inches(0.35), y + Inches(0.25), Inches(3), Inches(0.3),
             size=10, color=BRAND_TEAL, letter_spacing=3.0, bold=True)
    # Title
    add_text(s, title, x + Inches(0.35), y + Inches(0.6), quad_w - Inches(0.7), Inches(0.55),
             size=22, bold=True, color=TEXT_HEADING)
    # Protocols (large body)
    add_text(s, protocols, x + Inches(0.35), y + Inches(1.15), quad_w - Inches(0.7), Inches(0.55),
             size=15, color=TEXT_BODY, italic=False)

# Footer line: canonical vocabulary
add_text(s, "One canonical CNC vocabulary across every source. "
            "The same dashboard layout works across Fanuc, Brother, and Modbus-fronted machines.",
         Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.5),
         size=14, color=TEXT_CAPTION, italic=True, align="center")

# Coming-soon roadmap line
add_text(s, "COMING  ·  Siemens S7    ·    OPC UA Client    ·    HTTP / TCP sinks    ·    Linux host",
         Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.35),
         size=10, color=TEXT_MUTED, letter_spacing=2.5, align="center")

speaker_notes(s,
    "This is the trust-signal slide for industrial IT buyers in the room. "
    "If you're in front of a CNC-heavy crowd, spend time on the CNC column. "
    "If you're in front of a multi-vendor fleet, spend time on the canonical-vocabulary "
    "footer line — that's the unique differentiator.")

# ============================================================
# SLIDE 8 — Architecture at a glance
# ============================================================
s = add_blank(prs)
content_chrome(s, 8)

add_text(s, "ARCHITECTURE AT A GLANCE",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Edge collects  ·  Integration carries  ·  Intelligence aggregates  ·  Consumers consume",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.5),
         size=18, color=TEXT_BODY, italic=False, line_spacing=1.2)

# Embed architecture diagram (slide variant PNG, already 16:9)
arch_path = ASSETS_DIR / "architecture-diagram-v1-slide@2x.png"
if arch_path.exists():
    # The slide variant SVG is sized 3200x1800 (16:9). Embed into ~10.5x5.9".
    img_w = Inches(11.0)
    img_h = Inches(img_w.inches * (1800/3200))  # preserves 16:9
    img_x = (SLIDE_W - img_w) / 2
    img_y = Inches(1.65)
    s.shapes.add_picture(str(arch_path), img_x, img_y, img_w, img_h)
else:
    add_text(s, "[ architecture-diagram-v1-slide@2x.png missing ]",
             Inches(2), Inches(3), Inches(9), Inches(1), size=18, color=BRAND_TEAL, align="center")

# Caption below
add_text(s, "One EdgeConnect deploys at each plant. One EREMOS V2 tenant aggregates many sites. "
            "Standard MQTT and OPC UA make the integration interoperable with whatever else you run.",
         Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.65),
         size=13, color=TEXT_CAPTION, italic=True, align="center", line_spacing=1.3)

speaker_notes(s,
    "Walk left to right. One sentence per layer. Don't dwell on the OPC UA Server box unless "
    "asked — most plant managers don't care; most SCADA engineers do.")

# ============================================================
# SLIDE 9 — Why Elpis
# ============================================================
s = add_blank(prs)
content_chrome(s, 9)

add_text(s, "WHY ELPIS",
         Inches(0.7), Inches(0.5), Inches(6), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Differentiators, every one of them outcome-led.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.7),
         size=30, bold=True, color=TEXT_HEADING, line_spacing=1.1)

reasons = [
    ("New protocols ship without breaking the old ones",
     "Protocol-agnostic core by architecture, not by accident.", False),
    ("Built to run for years on a small box in a control cabinet",
     "Edge-first, not cloud-first. Store-and-forward by default.", False),
    ("Operators always know where the data flow broke",
     "Three-way diagnostics — source, pipeline, sink — by design.", False),
    ("Air-gapped factories are first-class",
     "RSA-signed JSON license, fully offline. No phone-home.", False),
    ("A lapsed license never stops production data",
     "Expiration blocks configuration changes only.", True),     # highlight
    ("AI proposes — humans decide",
     "Never silently alters the data path. Local-LLM-capable.", True),  # highlight
    ("Pay for the connectivity you actually use",
     "Per-edition packaging with modular per-protocol activation.", False),
    ("Built for industrial workloads",
     "Not adapted IoT software. OT-aware vocabulary throughout.", False),
]

# 4 rows × 2 columns
card_w = Inches(5.9)
card_h = Inches(1.05)
col_xs = [Inches(0.7), Inches(6.75)]
row_ys = [Inches(2.05), Inches(3.2), Inches(4.35), Inches(5.5)]

for i, (lead, sub, highlight) in enumerate(reasons):
    col = i % 2; row = i // 2
    x = col_xs[col]; y = row_ys[row]
    if highlight:
        # Highlighted row gets a more visible teal bar + slightly raised surface
        add_rect(s, x, y, card_w, card_h, fill_rgb=SURFACE_HERO, line_rgb=BRAND_TEAL)
        add_rect(s, x, y, Inches(0.08), card_h, fill_rgb=BRAND_TEAL)
        lead_color = TEXT_HEADING
    else:
        add_rect(s, x, y, Inches(0.04), card_h, fill_rgb=BORDER_STRONG)
        lead_color = TEXT_BODY
    add_text(s, lead, x + Inches(0.3), y + Inches(0.16), card_w - Inches(0.5), Inches(0.42),
             size=14, bold=True, color=lead_color, line_spacing=1.2)
    add_text(s, sub, x + Inches(0.3), y + Inches(0.58), card_w - Inches(0.5), Inches(0.4),
             size=12, color=TEXT_MUTED, line_spacing=1.3)

speaker_notes(s,
    "This is where you stop pitching features and start positioning category. "
    "Emphasize: 'AI proposes; humans decide' and 'A lapsed license never stops production "
    "data.' Both are credibility wins competitors can't replicate without rewriting their "
    "products.")

# ============================================================
# SLIDE 10 — Deploy incrementally
# ============================================================
s = add_blank(prs)
content_chrome(s, 10)

add_text(s, "DEPLOY INCREMENTALLY",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Start small. Expand without disruption.",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.9),
         size=34, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# Step arrow: 4 pills + arrow chain + "..." continuing
steps = ["One cell", "One line", "One plant", "Fleet"]
step_y = Inches(2.55)
step_w = Inches(2.2)
step_h = Inches(0.85)
gap = Inches(0.45)
total_w = step_w.inches * len(steps) + gap.inches * (len(steps) - 1) + 0.8  # +ellipsis
start_x = Inches((13.333 - total_w) / 2)

for i, label in enumerate(steps):
    x = start_x + Inches(i * (step_w.inches + gap.inches))
    is_last = (i == len(steps) - 1)
    if is_last:
        add_rect(s, x, step_y, step_w, step_h, fill_rgb=SURFACE_HERO, line_rgb=BRAND_TEAL)
    else:
        add_rect(s, x, step_y, step_w, step_h, fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE)
    add_text(s, label, x, step_y, step_w, step_h,
             size=18, bold=True, color=TEXT_HEADING, align="center", anchor="middle")
    # Arrow between steps
    if i < len(steps) - 1:
        ax = x + step_w
        ay = step_y + step_h / 2 - Inches(0.04)
        add_rect(s, ax + Inches(0.05), ay + Inches(0.03), gap - Inches(0.15), Inches(0.04),
                 fill_rgb=BRAND_TEAL)

# Ellipsis after last pill
ellipsis_x = start_x + Inches(len(steps) * (step_w.inches + gap.inches))
add_text(s, "...continues", ellipsis_x, step_y, Inches(1.5), step_h,
         size=14, color=TEXT_MUTED, italic=True, anchor="middle")

# Body paragraph
add_multiline(s, [
    {"text": "Start with one machine, one line, or one plant.",
     "size": 20, "color": TEXT_BODY, "bold": True},
    {"text": "EdgeConnect runs side-by-side with what you already have. "
             "EREMOS V2 onboards new sites without changing the platform underneath.",
     "size": 15, "color": TEXT_MUTED},
    {"text": "No big-bang cutover. No platform-wide upgrade that breaks the plants already running.",
     "size": 15, "color": TEXT_MUTED, "italic": True},
], Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.6),
   line_spacing=1.4, para_space_after=8, align="center")

# Callout box
co_y = Inches(5.85)
co_h = Inches(0.95)
add_rect(s, Inches(1.5), co_y, Inches(10.3), co_h,
         fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE)
add_rect(s, Inches(1.5), co_y, Inches(0.08), co_h, fill_rgb=BRAND_TEAL)
add_text(s, "Typical proof-of-value deployments begin with a single line or machine cell "
            "and expand incrementally once operationally validated.",
         Inches(1.75), co_y + Inches(0.15), Inches(10.0), co_h - Inches(0.3),
         size=14, color=TEXT_BODY, italic=True, anchor="middle", line_spacing=1.3)

speaker_notes(s,
    "This slide defuses the biggest unspoken objection: 'Will this disrupt my plant?' "
    "Read the callout box verbatim. Then pause.")

# ============================================================
# SLIDE 11 — Editions, modules, roadmap
# ============================================================
s = add_blank(prs)
content_chrome(s, 11)

add_text(s, "EDITIONS, MODULES, ROADMAP",
         Inches(0.7), Inches(0.5), Inches(8), Inches(0.35),
         size=11, color=TEXT_MUTED, letter_spacing=3.0, bold=True)

add_text(s, "Editions and modules",
         Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.65),
         size=28, bold=True, color=TEXT_HEADING, line_spacing=1.1)

# Three columns
ed_w = Inches(3.85)
ed_h = Inches(2.95)
ed_y = Inches(1.85)
ed_xs = [Inches(0.7), Inches(4.75), Inches(8.8)]

editions = [
    ("STARTER", "Single-plant deployments",
     ["Core EdgeConnect runtime", "Choice of 1 southbound module",
      "MQTT sink", "Diagnostics + store-and-forward"]),
    ("PROFESSIONAL", "Multi-protocol fleets",
     ["Everything in Starter", "Multiple southbound modules", "OPC UA Server",
      "Configurable routing + transforms"]),
    ("ENTERPRISE", "Multi-site + audit-grade",
     ["Everything in Professional", "Fleet-wide management",
      "Hash-chained config history", "Priority support + SLA"]),
]

for i, (name, sub, modules) in enumerate(editions):
    x = ed_xs[i]
    is_pro = (i == 1)  # highlight middle tier
    add_rect(s, x, ed_y, ed_w, ed_h,
             fill_rgb=SURFACE_HERO,
             line_rgb=BRAND_TEAL if is_pro else BORDER_SUBTLE)
    if is_pro:
        # Recommended ribbon
        add_rect(s, x + Inches(0.4), ed_y - Inches(0.12), Inches(1.5), Inches(0.28),
                 fill_rgb=BRAND_TEAL)
        add_text(s, "RECOMMENDED", x + Inches(0.4), ed_y - Inches(0.12), Inches(1.5), Inches(0.28),
                 size=9, bold=True, color=BG_DEEP, align="center", anchor="middle",
                 letter_spacing=2.0)
    add_text(s, name, x + Inches(0.3), ed_y + Inches(0.25), ed_w - Inches(0.5), Inches(0.4),
             size=15, bold=True, color=BRAND_TEAL, letter_spacing=2.5)
    add_text(s, sub, x + Inches(0.3), ed_y + Inches(0.68), ed_w - Inches(0.5), Inches(0.4),
             size=13, color=TEXT_MUTED, italic=True)
    add_rect(s, x + Inches(0.3), ed_y + Inches(1.15), ed_w - Inches(0.6), Inches(0.012),
             fill_rgb=BORDER_SUBTLE)
    for j, m in enumerate(modules):
        my = ed_y + Inches(1.3 + j * 0.36)
        add_text(s, "✓", x + Inches(0.3), my, Inches(0.3), Inches(0.3),
                 size=13, bold=True, color=BRAND_TEAL)
        add_text(s, m, x + Inches(0.55), my, ed_w - Inches(0.75), Inches(0.3),
                 size=12, color=TEXT_BODY)

# Module list footer
add_text(s, "Connectivity modules  ·  FOCAS2  ·  MT-LINKi  ·  MTConnect  ·  Brother HTTP  "
            "·  Modbus TCP  ·  MQTT  ·  OPC UA Server",
         Inches(0.7), Inches(4.95), Inches(11.9), Inches(0.4),
         size=11, color=TEXT_MUTED, letter_spacing=1.5, align="center", italic=True)

# Roadmap divider + second headline
add_rect(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.012), fill_rgb=BORDER_SUBTLE)

add_text(s, "ON THE ROADMAP",
         Inches(0.7), Inches(5.65), Inches(8), Inches(0.3),
         size=11, color=BRAND_TEAL, letter_spacing=3.0, bold=True)

# Roadmap items as pills
roadmap = [
    "OPC UA Client (southbound)",
    "Siemens S7 (southbound)",
    "HTTP / TCP sinks (northbound)",
    "Linux host support",
    "AI-assisted operations agents",
]
pill_y = Inches(6.05)
pill_h = Inches(0.5)
gap = Inches(0.15)
# Compute widths to fit
total_pill_text = sum(len(p) for p in roadmap)
available = 13.333 - 1.4  # 11.93
# Approximate per-pill width based on text length
pill_widths = [Inches(0.42 + 0.083 * len(p)) for p in roadmap]
sum_w = sum(w.inches for w in pill_widths) + gap.inches * (len(roadmap) - 1)
start_x = Inches((13.333 - sum_w) / 2)
cur_x = start_x
for i, label in enumerate(roadmap):
    w = pill_widths[i]
    add_rect(s, cur_x, pill_y, w, pill_h, fill_rgb=SURFACE_HERO, line_rgb=BORDER_SUBTLE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, label, cur_x, pill_y, w, pill_h,
             size=11, color=TEXT_BODY, align="center", anchor="middle")
    cur_x += w + gap

speaker_notes(s,
    "Cover editions in 30 seconds. The detail belongs in the follow-up conversation, not the "
    "pitch. The roadmap is here as a trust signal: 'We have a plan, and it serves the "
    "architecture, not the other way around.'")

# ============================================================
# SLIDE 12 — Next step
# ============================================================
s = add_blank(prs)
title_chrome(s)

# Subtle teal divider
add_rect(s, Inches(6.45), Inches(1.85), Inches(0.4), Inches(0.04), fill_rgb=BRAND_TEAL)

add_text(s, "NEXT STEP",
         Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.35),
         size=12, color=TEXT_MUTED, letter_spacing=4.0, bold=True, align="center")

add_text(s, "Bring us a real plant. We'll scope a real proof.",
         Inches(0.7), Inches(2.1), Inches(11.9), Inches(1.0),
         size=42, bold=True, color=TEXT_HEADING, align="center", line_spacing=1.1)

add_multiline(s, [
    {"text": "Bring us a representative plant — a controller mix, a target broker, an OEE definition.",
     "size": 18, "color": TEXT_BODY},
    {"text": "We will scope a proof of value against it.",
     "size": 18, "color": TEXT_BODY},
    {"text": "Demos run on real protocols against your real signals. Not canned data.",
     "size": 18, "color": TEXT_MUTED, "italic": True},
], Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.7),
   line_spacing=1.6, para_space_after=8, align="center")

# CTA pill
cta_w = Inches(3.5)
cta_h = Inches(0.65)
cta_x = (SLIDE_W - cta_w) / 2
cta_y = Inches(5.45)
add_rect(s, cta_x, cta_y, cta_w, cta_h, fill_rgb=BRAND_TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, "Book a scoping call",
         cta_x, cta_y, cta_w, cta_h,
         size=16, bold=True, color=BG_DEEP, align="center", anchor="middle", letter_spacing=1.5)

# Contact info block bottom
add_multiline(s, [
    {"text": "[ Presenter Name ]   ·   [ presenter@elpisitsolutions.com ]",
     "size": 13, "color": TEXT_BODY},
    {"text": "Elpis IT Solutions   ·   [ website ]   ·   [ phone ]",
     "size": 13, "color": TEXT_MUTED, "italic": True},
], Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.7),
   line_spacing=1.6, align="center")

# Footer brand mark
add_text(s, "ELPIS  ·  THINK · CREATE · ENABLE",
         Inches(0.7), Inches(7.15), Inches(11.9), Inches(0.3),
         size=8, color=TEXT_MUTED, align="center", letter_spacing=3.0)

speaker_notes(s,
    "Don't end with 'Any questions?' End with 'What would the first machine look like?' "
    "That question shifts the conversation from evaluating the pitch to evaluating their "
    "own plant.")

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------

prs.save(OUT_PPTX)
print(f"Wrote {OUT_PPTX}  ({OUT_PPTX.stat().st_size:,} bytes, {len(prs.slides)} slides)")
